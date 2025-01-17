#!/usr/bin/env python3
"""
Generate a PPTX presentation guiding Tamil speakers through
Japanese Hiragana and Katakana, including dakuten (が, ぎ, ぐ...) 
and handakuten (ぱ, ぴ, ぷ...). The script reads Romaji/Tamil
mappings from a CSV file, so changes in the CSV reflect in output.

We use 'Option 2': Hardcode the slide width & height in inches
(13.333" x 7.5") for 16:9, instead of retrieving them from
the Presentation object.
"""

import csv
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.enum.text import MSO_VERTICAL_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE

################################################################
# Hard-coded slide dimensions in inches (16:9 ratio)
################################################################
SLIDE_WIDTH_INCHES = 13.333
SLIDE_HEIGHT_INCHES = 7.5

################################################################
# 1) Basic (Non-Dakuten) Tables
################################################################
HIRAGANA_TABLE = [
  ["あ", "か", "さ", "た", "な", "は", "ま", "や", "ら", "わ"],
  ["い", "き", "し", "ち", "に", "ひ", "み", "",   "り", "" ],
  ["う", "く", "す", "つ", "ぬ", "ふ", "む", "ゆ", "る", "" ],
  ["え", "け", "せ", "て", "ね", "へ", "め", "",   "れ", "" ],
  ["お", "こ", "そ", "と", "の", "ほ", "も", "よ", "ろ", "を"],
]
HIRAGANA_N = "ん"

KATAKANA_TABLE = [
  ["ア", "カ", "サ", "タ", "ナ", "ハ", "マ", "ヤ", "ラ", "ワ"],
  ["イ", "キ", "シ", "チ", "ニ", "ヒ", "ミ", "",   "リ", "" ],
  ["ウ", "ク", "ス", "ツ", "ヌ", "フ", "ム", "ユ", "ル", "" ],
  ["エ", "ケ", "セ", "テ", "ネ", "ヘ", "メ", "",   "レ", "" ],
  ["オ", "コ", "ソ", "ト", "ノ", "ホ", "モ", "ヨ", "ロ", "ヲ"],
]
KATAKANA_N = "ン"

################################################################
# 2) Dakuten/Handakuten Tables
################################################################
HIRAGANA_DAKUTEN_TABLE = [
  ["が", "ぎ", "ぐ", "げ", "ご"],
  ["ざ", "じ", "ず", "ぜ", "ぞ"],
  ["だ", "ぢ", "づ", "で", "ど"],
  ["ば", "び", "ぶ", "べ", "ぼ"],
  ["ぱ", "ぴ", "ぷ", "ぺ", "ぽ"],  # handakuten row
]

KATAKANA_DAKUTEN_TABLE = [
  ["ガ", "ギ", "グ", "ゲ", "ゴ"],
  ["ザ", "ジ", "ズ", "ゼ", "ゾ"],
  ["ダ", "ヂ", "ヅ", "デ", "ド"],
  ["バ", "ビ", "ブ", "ベ", "ボ"],
  ["パ", "ピ", "プ", "ペ", "ポ"],
]

################################################################
# 3) We'll create a dictionary at runtime from CSV
################################################################
ROMAJI_TAMIL_MAP = {}

def load_mapping_from_csv(csv_path):
    """
    Reads a CSV with columns: hiragana, katakana, tamil, romaji
    and populates the global ROMAJI_TAMIL_MAP so that:
      ROMAJI_TAMIL_MAP[hira] = (romaji, tamil)
      ROMAJI_TAMIL_MAP[kata] = (romaji, tamil)
    """
    with open(csv_path, mode="r", encoding="utf-8") as f:
        reader = csv.DictReader(f)
        for row in reader:
            hira = row["hiragana"].strip()
            kata = row["katakana"].strip()
            tam  = row["tamil"].strip()
            rom  = row["romaji"].strip()

            # For each line, store for both hiragana and katakana
            ROMAJI_TAMIL_MAP[hira] = (rom, tam)
            ROMAJI_TAMIL_MAP[kata] = (rom, tam)

################################################################
# 4) Series for Focus Slides
################################################################
GOJUON_SERIES = [
    ("A Series",  ["あ", "い", "う", "え", "お"],  ["ア", "イ", "ウ", "エ", "オ"]),
    ("Ka Series", ["か", "き", "く", "け", "こ"],  ["カ", "キ", "ク", "ケ", "コ"]),
    ("Sa Series", ["さ", "し", "す", "せ", "そ"],  ["サ", "シ", "ス", "セ", "ソ"]),
    ("Ta Series", ["た", "ち", "つ", "て", "と"],  ["タ", "チ", "ツ", "テ", "ト"]),
    ("Na Series", ["な", "に", "ぬ", "ね", "の"],  ["ナ", "ニ", "ヌ", "ネ", "ノ"]),
    ("Ha Series", ["は", "ひ", "ふ", "へ", "ほ"],  ["ハ", "ヒ", "フ", "ヘ", "ホ"]),
    ("Ma Series", ["ま", "み", "む", "め", "も"],  ["マ", "ミ", "ム", "メ", "モ"]),
    ("Ya Series", ["や", "ゆ", "よ"],             ["ヤ", "ユ", "ヨ"]),
    ("Ra Series", ["ら", "り", "る", "れ", "ろ"],  ["ラ", "リ", "ル", "レ", "ロ"]),
    ("Wa/N Series",["わ", "を", "ん"],            ["ワ", "ヲ", "ン"]),
]

DAKUTEN_SERIES = [
    ("Ga Series", ["が", "ぎ", "ぐ", "げ", "ご"], ["ガ", "ギ", "グ", "ゲ", "ゴ"]),
    ("Za Series", ["ざ", "じ", "ず", "ぜ", "ぞ"], ["ザ", "ジ", "ズ", "ゼ", "ゾ"]),
    ("Da Series", ["だ", "ぢ", "づ", "で", "ど"], ["ダ", "ヂ", "ヅ", "デ", "ド"]),
    ("Ba Series", ["ば", "び", "ぶ", "べ", "ぼ"], ["バ", "ビ", "ブ", "ベ", "ボ"]),
    ("Pa Series", ["ぱ", "ぴ", "ぷ", "ぺ", "ぽ"], ["パ", "ピ", "プ", "ペ", "ポ"]),
]

################################################################
# 5) Helper Functions
################################################################
def create_table_for_syllabary(
    slide,
    table_data,
    slide_title,
    top=1.0,
    left=0.5,
    col_width=1.0,
    row_height=1.2,
    font_main=30,
    font_sub=14
):
    """
    Creates a table on the given slide (chart style).
    Each cell shows:
      - Big script char on top
      - Smaller line with (Romaji | Tamil) below

    We do NOT retrieve slide dimensions from 'slide' or 'prs'.
    Instead, we rely on SLIDE_WIDTH_INCHES / SLIDE_HEIGHT_INCHES
    for simple auto-scaling.

    If the user modifies or extends the CSV, the dictionary
    will reflect those changes automatically, so new lines
    appear or missing lines show blank.
    """
    slide_width_in = SLIDE_WIDTH_INCHES
    slide_height_in = SLIDE_HEIGHT_INCHES

    rows = len(table_data)
    cols = max(len(r) for r in table_data) if rows > 0 else 0

    # Title
    title_shape = slide.shapes.add_textbox(
        Inches(left),
        Inches(top),
        Inches(col_width * cols),
        Inches(0.5)
    )
    tf = title_shape.text_frame
    tf.text = slide_title
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    tf.paragraphs[0].runs[0].font.size = Pt(font_main + 6)
    tf.paragraphs[0].runs[0].font.bold = True

    # Table position below the title
    table_top = top + 0.7
    table_left = left
    table_width = col_width * cols
    table_height = row_height * rows

    # Simple scaling if table too wide:
    margin_w = 1.0  
    available_w = slide_width_in - margin_w
    if table_width > available_w:
        scale_factor = available_w / table_width
        col_width *= scale_factor
        table_width = col_width * cols

    # Similarly for height
    margin_h = 1.5  
    available_h = slide_height_in - margin_h
    if table_height > available_h:
        scale_factor_h = available_h / table_height
        row_height *= scale_factor_h
        table_height = row_height * rows

    # Create the table
    graphic_frame = slide.shapes.add_table(
        rows, cols,
        Inches(table_left),
        Inches(table_top),
        Inches(table_width),
        Inches(table_height)
    )
    table = graphic_frame.table

    # Set uniform column widths & row heights
    for col_idx in range(cols):
        table.columns[col_idx].width = Inches(col_width)
    for row_idx in range(rows):
        table.rows[row_idx].height = Inches(row_height)

    # Fill each cell
    for r_idx, row_data in enumerate(table_data):
        for c_idx, char in enumerate(row_data):
            cell = table.cell(r_idx, c_idx)
            if not char:
                cell.text = ""
                continue

            # Pull romaji/tamil from the dictionary
            romaji, tamil = ROMAJI_TAMIL_MAP.get(char, ("", ""))

            cell.text_frame.clear()
            cell.text_frame.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE

            # Big script char
            p1 = cell.text_frame.add_paragraph()
            p1.text = char
            p1.alignment = PP_ALIGN.CENTER
            run1 = p1.runs[0]
            run1.font.size = Pt(font_main)
            run1.font.bold = True

            # Smaller line: Romaji | Tamil
            p2 = cell.text_frame.add_paragraph()
            p2.text = f"{romaji}  |  {tamil}"
            p2.alignment = PP_ALIGN.CENTER
            run2 = p2.runs[0]
            run2.font.size = Pt(font_sub)
            run2.font.bold = False

def create_centered_textbox(
    slide,
    text,
    left,
    top,
    width,
    height,
    font_size=60,
    bold=True,
    alignment=PP_ALIGN.CENTER
):
    """
    Creates a textbox shape on the slide, centered horizontally & vertically,
    with the given font size & boldness.
    """
    textbox = slide.shapes.add_textbox(
        Inches(left), Inches(top),
        Inches(width), Inches(height)
    )
    text_frame = textbox.text_frame
    text_frame.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE

    p = text_frame.paragraphs[0]
    p.alignment = alignment

    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = RGBColor(0, 0, 0)  # black text

    return textbox

################################################################
# 6) Main Generation Logic
################################################################
def main():
    # 1) Load the CSV for mappings
    load_mapping_from_csv("mapping.csv")  # Adjust path as needed

    # 2) Setup the Presentation
    prs = Presentation()
    prs.slide_width = Inches(SLIDE_WIDTH_INCHES)
    prs.slide_height = Inches(SLIDE_HEIGHT_INCHES)

    blank_layout = prs.slide_layouts[6]  # blank layout

    # ----------------------------------------------------------------
    # A) Basic Hiragana Overview
    # ----------------------------------------------------------------
    slide1 = prs.slides.add_slide(blank_layout)
    create_table_for_syllabary(
        slide1,
        HIRAGANA_TABLE,
        "Hiragana (ひらがな) Overview",
        top=0.5, left=0.5,
        col_width=1.1, row_height=1.0,
        font_main=32, font_sub=14
    )
    # Add "ん" in a separate textbox if desired
    create_centered_textbox(
        slide1,
        f"{HIRAGANA_N}\nN | {ROMAJI_TAMIL_MAP.get(HIRAGANA_N, ('',''))[1]}",  
        left=11.0, top=5.5, width=1.2, height=1.0,
        font_size=24, bold=False
    )

    # ----------------------------------------------------------------
    # B) Basic Katakana Overview
    # ----------------------------------------------------------------
    slide2 = prs.slides.add_slide(blank_layout)
    create_table_for_syllabary(
        slide2,
        KATAKANA_TABLE,
        "Katakana (カタカナ) Overview",
        top=0.5, left=0.5,
        col_width=1.1, row_height=1.0,
        font_main=32, font_sub=14
    )
    create_centered_textbox(
        slide2,
        f"{KATAKANA_N}\nN | {ROMAJI_TAMIL_MAP.get(KATAKANA_N, ('',''))[1]}",
        left=11.0, top=5.5, width=1.2, height=1.0,
        font_size=24, bold=False
    )

    # ----------------------------------------------------------------
    # C) Hiragana Dakuten/Handakuten
    # ----------------------------------------------------------------
    slide3 = prs.slides.add_slide(blank_layout)
    create_table_for_syllabary(
        slide3,
        HIRAGANA_DAKUTEN_TABLE,
        "Hiragana Dakuten/Handakuten",
        top=0.5, left=0.7,
        col_width=1.1, row_height=1.0,
        font_main=32, font_sub=14
    )

    # ----------------------------------------------------------------
    # D) Katakana Dakuten/Handakuten
    # ----------------------------------------------------------------
    slide4 = prs.slides.add_slide(blank_layout)
    create_table_for_syllabary(
        slide4,
        KATAKANA_DAKUTEN_TABLE,
        "Katakana Dakuten/Handakuten",
        top=0.5, left=0.7,
        col_width=1.1, row_height=1.0,
        font_main=32, font_sub=14
    )

    # ----------------------------------------------------------------
    # E) Focus Slides for Basic Gojuon
    # ----------------------------------------------------------------
    for (series_name, hira_list, kata_list) in GOJUON_SERIES:
        # E1) Series Overview Slide
        overview_slide = prs.slides.add_slide(blank_layout)
        create_centered_textbox(
            overview_slide,
            series_name,
            left=2.0, top=1.0,
            width=8.0, height=1.0,
            font_size=60, bold=True
        )
        # Show the row
        hira_str = " ".join(hira_list)
        kata_str = " ".join(kata_list)
        create_centered_textbox(
            overview_slide,
            f"Hiragana: {hira_str}\nKatakana: {kata_str}",
            left=2.0, top=3.0,
            width=8.0, height=2.0,
            font_size=36, bold=False
        )

        # E2) Individual Focus Slides
        for i in range(len(hira_list)):
            h = hira_list[i]
            k = kata_list[i]

            # For each char, get the romaji/tamil from the dictionary
            (romaji_h, tamil_h) = ROMAJI_TAMIL_MAP.get(h, ("", ""))
            (romaji_k, tamil_k) = ROMAJI_TAMIL_MAP.get(k, ("", ""))

            focus_slide = prs.slides.add_slide(blank_layout)
            # Large side-by-side
            create_centered_textbox(
                focus_slide,
                f"{h}    {k}",
                left=3.0, top=2.0,
                width=7.0, height=1.5,
                font_size=120, bold=True
            )
            # Smaller Romaji | Tamil (from Hiragana or unify them)
            create_centered_textbox(
                focus_slide,
                f"{romaji_h} | {tamil_h}",
                left=3.0, top=4.0,
                width=7.0, height=1.0,
                font_size=50, bold=False
            )

    # ----------------------------------------------------------------
    # F) Focus Slides for Dakuten Series
    # ----------------------------------------------------------------
    for (series_name, hira_list, kata_list) in DAKUTEN_SERIES:
        # F1) Series Overview
        overview_slide = prs.slides.add_slide(blank_layout)
        create_centered_textbox(
            overview_slide,
            series_name,
            left=2.0, top=1.0,
            width=8.0, height=1.0,
            font_size=60, bold=True
        )
        hira_str = " ".join(hira_list)
        kata_str = " ".join(kata_list)
        create_centered_textbox(
            overview_slide,
            f"Hiragana: {hira_str}\nKatakana: {kata_str}",
            left=2.0, top=3.0,
            width=8.0, height=2.0,
            font_size=36, bold=False
        )

        # F2) Individual Focus Slides
        for i in range(len(hira_list)):
            h = hira_list[i]
            k = kata_list[i]

            (romaji_h, tamil_h) = ROMAJI_TAMIL_MAP.get(h, ("", ""))
            (romaji_k, tamil_k) = ROMAJI_TAMIL_MAP.get(k, ("", ""))

            focus_slide = prs.slides.add_slide(blank_layout)
            # Big top
            create_centered_textbox(
                focus_slide,
                f"{h}    {k}",
                left=3.0, top=2.0,
                width=7.0, height=1.5,
                font_size=120, bold=True
            )
            # Smaller bottom
            create_centered_textbox(
                focus_slide,
                f"{romaji_h} | {tamil_h}",
                left=3.0, top=4.0,
                width=7.0, height=1.0,
                font_size=50, bold=False
            )

    # ----------------------------------------------------------------
    # G) Save
    # ----------------------------------------------------------------
    output_name = "Japanese_Guide_for_Tamil_Speakers_v4.pptx"
    prs.save(output_name)
    print(f"Presentation saved as '{output_name}'.")

if __name__ == "__main__":
    main()
